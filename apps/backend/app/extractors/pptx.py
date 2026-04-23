"""PPTX extraction via python-pptx — slide-by-slide text dump."""

from __future__ import annotations

from pathlib import Path

from .base import ExtractionError, ExtractResult


class PptxExtractor:
    extensions = ("pptx",)

    def extract(self, path: Path) -> ExtractResult:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError as e:  # pragma: no cover
            raise ExtractionError("python-pptx is not installed") from e

        try:
            prs = Presentation(str(path))
        except Exception as e:
            raise ExtractionError(f"Failed to open PPTX {path.name}: {e}") from e

        parts: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                # Each shape may or may not expose text — guard accordingly
                text = getattr(shape, "text", None)
                if text:
                    texts.append(text)

                # Tables aren't exposed via shape.text in older versions
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text_frame.text if cell.text_frame else ""
                            for cell in row.cells
                        )
                        texts.append(row_text)

            # Include slide notes — often where speaker context lives
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    texts.append(f"[Notes] {notes}")

            if texts:
                parts.append(f"=== Slide {idx} ===\n" + "\n".join(texts))
            else:
                parts.append(f"=== Slide {idx} ===\n(no text)")

        return ExtractResult(
            text="\n\n".join(parts),
            metadata={"slides": len(prs.slides)},
        )
